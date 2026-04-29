"""
gemma4_extractor.py

Extracts content from documents using a local Gemma4 model via Ollama.
Unified approach: everything → PDF → images → Gemma4 vision model.

Supports: PDF, images (jpg/png/gif/webp), DOCX, PPTX, XLSX, CSV, TXT, Markdown, JSON, XML, HTML
Requires: Ollama running on localhost:11434 with gemma4 model pulled.
         pip install ollama pymupdf pdf2image
"""

from __future__ import annotations

import base64
import io
from pathlib import Path

METHOD = "gemma4"

_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
_TEXT_EXTS = {".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml", ".yaml", ".yml"}
_CONVERT_EXTS = {".docx", ".pptx", ".xlsx"} | _TEXT_EXTS

_EXTRACTION_PROMPT = """Extract all content from this document and return it as Markdown.

Rules:
- Preserve the exact order content appears in the document — top to bottom, do not regroup.
- Use ## for section headings and ### for sub-headings, matching the document hierarchy.
- Render label:value pairs as **Label:** Value on their own line.
- Render tables as Markdown tables with headers and all rows.
- Copy all text verbatim — do not rephrase, summarize, or add commentary.
- Do not add any text that is not in the document.
- For charts, graphs, or plots that contain legible numeric or categorical data, render the
  data as a Markdown table (rows × columns matching the chart axes and legend), then
  immediately follow it with a one-line figure note describing the chart type, title, and
  any annotation not captured in the table:
  > **[Figure]** <chart type, title, axis labels, color encoding, key annotations>
- For every other visual element (logo, photograph, diagram, sketch, signature, illustration,
  stamp, or any chart whose data values cannot be read) insert only a figure block:
  > **[Figure]** <description of shape, colors, symbols, text, spatial layout, and meaning>
  Do not just name it — describe it as if explaining to someone who cannot see it.
"""

_MODEL = "gemma4:latest"  # Smaller, faster (9.6GB)
_MAX_TOKENS = 4000  # Budget for thinking + output


def _html_to_pdf_bytes(html: str) -> bytes:
    """Convert HTML to PDF bytes using fitz.Story."""
    import fitz

    story = fitz.Story(html)
    mediabox = fitz.paper_rect("a4")
    where = mediabox + (50, 50, -50, -50)

    buf = io.BytesIO()
    writer = fitz.DocumentWriter(buf)
    more = True
    while more:
        dev = writer.begin_page(mediabox)
        more, _ = story.place(where)
        story.draw(dev)
        writer.end_page()
    writer.close()
    return buf.getvalue()


def _docx_to_html(content: bytes) -> str:
    """Extract DOCX as HTML with inline base64 images."""
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(content))

    img_map: dict[str, str] = {}
    for rId, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            b64 = base64.b64encode(blob).decode()
            img_map[rId] = f"data:image/png;base64,{b64}"

    parts = ["<html><body>"]

    for block in doc.element.body:
        tag = block.tag.split("}")[-1]

        if tag == "p":
            style = ""
            if hasattr(block, "style") and block.style and hasattr(block.style, "name"):
                style = block.style.name
            text = block.text.strip() if hasattr(block, "text") else ""

            if not text:
                parts.append("<br/>")
                continue

            if "Heading 1" in style:
                parts.append(f"<h1>{text}</h1>")
            elif "Heading 2" in style:
                parts.append(f"<h2>{text}</h2>")
            elif "Heading 3" in style:
                parts.append(f"<h3>{text}</h3>")
            elif "List" in style:
                parts.append(f"<li>{text}</li>")
            else:
                parts.append(f"<p>{text}</p>")

            for blip in block.iter(qn("a:blip")):
                embed = blip.get(qn("r:embed"))
                if embed and embed in img_map:
                    parts.append(f'<img src="{img_map[embed]}" style="max-width:100%;"/>')

        elif tag == "tbl":
            parts.append("<table border='1' cellpadding='4'>")
            for row in block:
                if row.tag.split("}")[-1] != "tr":
                    continue
                parts.append("<tr>")
                for cell in row:
                    if cell.tag.split("}")[-1] not in ("tc",):
                        continue
                    cell_text = "".join(p.text or "" for p in cell.iter() if p.tag.endswith("}t"))
                    parts.append(f"<td>{cell_text}</td>")
                parts.append("</tr>")
            parts.append("</table>")

    parts.append("</body></html>")
    return "\n".join(parts)


def _pptx_to_html(content: bytes) -> str:
    """Extract PPTX as HTML with inline base64 images."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(content))
    parts = ["<html><body>"]

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts.append(f"<h2>Slide {slide_num}</h2>")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip() if para.text else ""
                    if text:
                        parts.append(f"<p>{text}</p>")

            if shape.has_table:
                parts.append("<table border='1' cellpadding='4'>")
                for row in shape.table.rows:
                    parts.append("<tr>")
                    for cell in row.cells:
                        parts.append(f"<td>{cell.text.strip()}</td>")
                    parts.append("</tr>")
                parts.append("</table>")

            if shape.has_chart:
                chart = shape.chart
                parts.append("<pre>")
                try:
                    for series in chart.series:
                        vals = list(series.values)
                        parts.append(f"Series '{series.name}': {vals}")
                except Exception:
                    parts.append("[chart — data unavailable]")
                parts.append("</pre>")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    b64 = base64.b64encode(blob).decode()
                    parts.append(f'<img src="data:image/png;base64,{b64}" style="max-width:100%;"/>')
                except Exception:
                    pass

        parts.append("<hr/>")

    parts.append("</body></html>")
    return "\n".join(parts)


def _xlsx_to_html(content: bytes) -> str:
    """Extract XLSX as HTML."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = ["<html><body>"]
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"<h2>{name}</h2><table border='1' cellpadding='4'>")
        for row in ws.iter_rows(values_only=True):
            parts.append("<tr>" + "".join(f"<td>{v or ''}</td>" for v in row) + "</tr>")
        parts.append("</table>")
    parts.append("</body></html>")
    return "\n".join(parts)


def _pdf_to_images(content: bytes) -> list[bytes]:
    """Convert PDF to list of page images (PNG bytes)."""
    from pdf2image import convert_from_bytes

    images = convert_from_bytes(content, fmt="png")
    png_bytes_list = []
    for img in images:
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        png_bytes_list.append(buf.getvalue())
    return png_bytes_list


def _call_with_image(client, b64_image: str, file_name: str) -> str:
    """Send image to Gemma4 with optimal settings for OCR/document parsing."""
    prompt = f"{_EXTRACTION_PROMPT}\n\nDOCUMENT: {file_name}"

    response = client.generate(
        model=_MODEL,
        prompt=prompt,
        images=[b64_image],
        stream=False,
        options={
            "temperature": 1.0,      # Gemma 4 recommended
            "top_p": 0.95,           # Gemma 4 recommended
            "top_k": 64,             # Gemma 4 recommended
            "num_predict": _MAX_TOKENS,  # Budget for thinking + output
        },
    )

    return response.get("response", "")


def extract(file_path: str) -> dict:
    """
    Unified extraction: everything → PDF → images → Gemma4 vision model.
    Same approach as Claude extractor but using local Gemma4.
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    content = path.read_bytes()
    warnings_list = []

    try:
        import ollama
    except ImportError:
        return {
            "file": path.name,
            "method": METHOD,
            "raw_markdown": "",
            "raw_text_chars": 0,
            "warnings": ["ollama package not installed"],
        }

    client = ollama

    try:
        # Step 1a: Handle native images (send directly, no PDF conversion)
        if ext in _IMAGE_EXTS:
            b64_image = base64.standard_b64encode(content).decode()
            raw_markdown = _call_with_image(client, b64_image, path.name)
        else:
            # Step 1b: Convert everything else to PDF
            if ext == ".pdf":
                pdf_bytes = content
            elif ext == ".docx":
                html = _docx_to_html(content)
                pdf_bytes = _html_to_pdf_bytes(html)
            elif ext == ".pptx":
                html = _pptx_to_html(content)
                pdf_bytes = _html_to_pdf_bytes(html)
            elif ext == ".xlsx":
                html = _xlsx_to_html(content)
                pdf_bytes = _html_to_pdf_bytes(html)
            elif ext in _TEXT_EXTS:
                text = content.decode("utf-8", errors="replace")
                html = f"<html><body><pre>{text}</pre></body></html>"
                pdf_bytes = _html_to_pdf_bytes(html)
            else:
                try:
                    text = content.decode("utf-8", errors="replace")
                    html = f"<html><body><pre>{text}</pre></body></html>"
                    pdf_bytes = _html_to_pdf_bytes(html)
                    warnings_list.append(f"Unknown type {ext}; treated as plain text")
                except Exception:
                    return {
                        "file": path.name,
                        "method": METHOD,
                        "raw_markdown": "",
                        "raw_text_chars": 0,
                        "warnings": [f"Unsupported file type: {ext}"],
                    }

            # Step 2: Convert PDF to images (Ollama vision needs image format, not raw PDF)
            images = _pdf_to_images(pdf_bytes)

            # Step 3: Send each page image to Gemma4
            responses = []
            for i, img_bytes in enumerate(images, start=1):
                b64_image = base64.standard_b64encode(img_bytes).decode()
                page_label = f"{path.stem} (page {i})" if len(images) > 1 else path.stem
                resp = _call_with_image(client, b64_image, page_label)
                if resp:
                    if len(images) > 1:
                        responses.append(f"--- Page {i} ---\n{resp}")
                    else:
                        responses.append(resp)

            raw_markdown = "\n".join(responses)

    except Exception as exc:
        error_msg = f"Gemma4 extraction error: {exc}"
        if "Connection refused" in str(exc) or "not installed" in str(exc):
            error_msg = f"Ollama not running or gemma4 model unavailable: {exc}"
        return {
            "file": path.name,
            "method": METHOD,
            "raw_markdown": "",
            "raw_text_chars": 0,
            "warnings": [error_msg],
        }

    return {
        "file": path.name,
        "method": METHOD,
        "raw_text_chars": len(raw_markdown),
        "raw_markdown": raw_markdown,
        "warnings": warnings_list,
    }
