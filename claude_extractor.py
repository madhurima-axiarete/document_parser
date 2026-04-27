"""
claude_extractor.py

Extracts content from documents using the Anthropic Claude API.
All non-native formats (DOCX, PPTX, XLSX, text) are converted to PDF via HTML
before sending to Claude, ensuring one unified code path.

Supports: PDF, images (jpg/png/gif/webp), DOCX, PPTX, XLSX, CSV, TXT, Markdown, JSON, XML, HTML
Requires: ANTHROPIC_API_KEY env var.
"""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

METHOD = "claude"

_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
_TEXT_EXTS = {".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml", ".yaml", ".yml"}
_CONVERT_EXTS = {".docx", ".pptx", ".xlsx"} | _TEXT_EXTS

_EXTRACTION_PROMPT = """Extract all content from this document and return it as Markdown.

Rules:
- Preserve the exact order content appears in the document — top to bottom, do not regroup.
- Use ## for section headings and ### for sub-headings, matching the document hierarchy.
- Render label:value pairs as **Label:** Value on their own line.
- Render tables as Markdown tables with headers and all rows.
- Copy all text verbatim — do not rephrase, summarize, or add commentary.
- Do not add any text that is not in the document.
- For charts, graphs, or plots that contain legible numeric or categorical data, render the
  data as a Markdown table (rows × columns matching the chart axes and legend), then
  immediately follow it with a one-line figure note describing the chart type, title, and
  any annotation not captured in the table:
  > **[Figure]** <chart type, title, axis labels, color encoding, key annotations>
- For every other visual element (logo, photograph, diagram, sketch, signature, illustration,
  stamp, or any chart whose data values cannot be read) insert only a figure block:
  > **[Figure]** <description of shape, colors, symbols, text, spatial layout, and meaning>
  Do not just name it — describe it as if explaining to someone who cannot see it.
"""

_MODEL = "claude-sonnet-4-6"
_MAX_TOKENS = 64000  # requires streaming for >16K


def _stream_text(client, messages: list) -> str:
    """Stream a response and return the full text (no output-token ceiling)."""
    with client.messages.stream(
        model=_MODEL,
        max_tokens=_MAX_TOKENS,
        messages=messages,
    ) as stream:
        return stream.get_final_text()


def _call_with_pdf(client, content: bytes, file_name: str) -> str:
    return _stream_text(client, [{
        "role": "user",
        "content": [
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": base64.standard_b64encode(content).decode(),
                },
            },
            {"type": "text", "text": f"{_EXTRACTION_PROMPT}\n\nDOCUMENT: {file_name}"},
        ],
    }])


def _call_with_image(client, content: bytes, ext: str, file_name: str) -> str:
    media_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
                   ".gif": "image/gif", ".webp": "image/webp"}
    return _stream_text(client, [{
        "role": "user",
        "content": [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_types.get(ext, "image/jpeg"),
                    "data": base64.standard_b64encode(content).decode(),
                },
            },
            {"type": "text", "text": f"{_EXTRACTION_PROMPT}\n\nDOCUMENT: {file_name}"},
        ],
    }])


def _sniff_image_media_type(blob: bytes) -> str:
    """Detect image format from magic bytes."""
    if blob[:2] == b'\xff\xd8':
        return "image/jpeg"
    if blob[:4] == b'\x89PNG':
        return "image/png"
    if blob[:4] in (b'GIF8', b'GIF9'):
        return "image/gif"
    return "image/png"


def _html_to_pdf_bytes(html: str) -> bytes:
    """Convert HTML to PDF bytes using fitz.Story."""
    import fitz

    story = fitz.Story(html)
    mediabox = fitz.paper_rect("a4")
    where = mediabox + (50, 50, -50, -50)

    buf = io.BytesIO()
    writer = fitz.DocumentWriter(buf)
    more = True
    while more:
        dev = writer.begin_page(mediabox)
        more, _ = story.place(where)
        story.draw(dev)
        writer.end_page()
    writer.close()
    return buf.getvalue()


def _docx_to_html(content: bytes) -> str:
    """Extract DOCX as HTML with inline base64 images."""
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(content))

    img_map: dict[str, str] = {}
    for rId, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            mt = _sniff_image_media_type(blob)
            img_map[rId] = f"data:{mt};base64,{base64.b64encode(blob).decode()}"

    parts = ["<html><body>"]

    for block in doc.element.body:
        tag = block.tag.split("}")[-1]

        if tag == "p":
            style = ""
            if hasattr(block, "style") and block.style and hasattr(block.style, "name"):
                style = block.style.name
            text = block.text.strip() if hasattr(block, "text") else ""

            if not text:
                parts.append("<br/>")
                continue

            if "Heading 1" in style:
                parts.append(f"<h1>{text}</h1>")
            elif "Heading 2" in style:
                parts.append(f"<h2>{text}</h2>")
            elif "Heading 3" in style:
                parts.append(f"<h3>{text}</h3>")
            elif "List" in style:
                parts.append(f"<li>{text}</li>")
            else:
                parts.append(f"<p>{text}</p>")

            # Check for inline images
            for blip in block.iter(qn("a:blip")):
                embed = blip.get(qn("r:embed"))
                if embed and embed in img_map:
                    parts.append(f'<img src="{img_map[embed]}" style="max-width:100%;"/>')

        elif tag == "tbl":
            parts.append("<table border='1' cellpadding='4'>")
            for row in block:
                if row.tag.split("}")[-1] != "tr":
                    continue
                parts.append("<tr>")
                for cell in row:
                    if cell.tag.split("}")[-1] not in ("tc",):
                        continue
                    cell_text = "".join(p.text or "" for p in cell.iter() if p.tag.endswith("}t"))
                    parts.append(f"<td>{cell_text}</td>")
                parts.append("</tr>")
            parts.append("</table>")

    parts.append("</body></html>")
    return "\n".join(parts)


def _pptx_to_html(content: bytes) -> str:
    """Extract PPTX as HTML with inline base64 images."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(content))
    parts = ["<html><body>"]

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts.append(f"<h2>Slide {slide_num}</h2>")

        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip() if para.text else ""
                    if text:
                        parts.append(f"<p>{text}</p>")

            # Tables
            if shape.has_table:
                parts.append("<table border='1' cellpadding='4'>")
                for row in shape.table.rows:
                    parts.append("<tr>")
                    for cell in row.cells:
                        parts.append(f"<td>{cell.text.strip()}</td>")
                    parts.append("</tr>")
                parts.append("</table>")

            # Charts
            if shape.has_chart:
                chart = shape.chart
                parts.append("<pre>")
                try:
                    for series in chart.series:
                        vals = list(series.values)
                        parts.append(f"Series '{series.name}': {vals}")
                except Exception:
                    parts.append("[chart — data unavailable]")
                parts.append("</pre>")

            # Embedded images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    mt = _sniff_image_media_type(blob)
                    b64 = base64.b64encode(blob).decode()
                    parts.append(f'<img src="data:{mt};base64,{b64}" style="max-width:100%;"/>')
                except Exception:
                    pass

        parts.append("<hr/>")

    parts.append("</body></html>")
    return "\n".join(parts)


def _xlsx_to_html(content: bytes) -> str:
    """Extract XLSX as HTML."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = ["<html><body>"]
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"<h2>{name}</h2><table border='1' cellpadding='4'>")
        for row in ws.iter_rows(values_only=True):
            parts.append("<tr>" + "".join(f"<td>{v or ''}</td>" for v in row) + "</tr>")
        parts.append("</table>")
    parts.append("</body></html>")
    return "\n".join(parts)


def extract(file_path: str) -> dict:
    """Send file to Claude and return the Markdown response."""
    path = Path(file_path)
    ext = path.suffix.lower()
    content = path.read_bytes()
    warnings_list = []

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {
            "file": path.name, "method": METHOD,
            "raw_markdown": "", "raw_text_chars": 0,
            "warnings": ["ANTHROPIC_API_KEY not set — skipping"],
        }

    try:
        import anthropic
        from anthropic._exceptions import RequestTooLargeError
        client = anthropic.Anthropic(api_key=api_key)
    except ImportError:
        return {
            "file": path.name, "method": METHOD,
            "raw_markdown": "", "raw_text_chars": 0,
            "warnings": ["anthropic package not installed"],
        }

    try:
        if ext in _IMAGE_EXTS:
            raw_markdown = _call_with_image(client, content, ext, path.name)
        elif ext == ".pdf":
            raw_markdown = _call_with_pdf(client, content, path.name)
        elif ext in _CONVERT_EXTS:
            if ext == ".docx":
                html = _docx_to_html(content)
            elif ext == ".pptx":
                html = _pptx_to_html(content)
            elif ext == ".xlsx":
                html = _xlsx_to_html(content)
            else:
                text = content.decode("utf-8", errors="replace")
                html = f"<html><body><pre>{text}</pre></body></html>"
            pdf_bytes = _html_to_pdf_bytes(html)
            raw_markdown = _call_with_pdf(client, pdf_bytes, path.name)
        else:
            try:
                text = content.decode("utf-8", errors="replace")
                html = f"<html><body><pre>{text}</pre></body></html>"
                pdf_bytes = _html_to_pdf_bytes(html)
                raw_markdown = _call_with_pdf(client, pdf_bytes, path.name)
                warnings_list.append(f"Unknown type {ext}; treated as plain text")
            except Exception:
                return {
                    "file": path.name, "method": METHOD,
                    "raw_markdown": "", "raw_text_chars": 0,
                    "warnings": [f"Unsupported file type: {ext}"],
                }

    except RequestTooLargeError:
        mb = len(content) / (1024 * 1024)
        return {
            "file": path.name, "method": METHOD,
            "raw_markdown": "", "raw_text_chars": 0,
            "warnings": [f"File too large for Claude API ({mb:.1f} MB); skipped"],
        }
    except Exception as exc:
        return {
            "file": path.name, "method": METHOD,
            "raw_markdown": "", "raw_text_chars": 0,
            "warnings": [f"Claude API error: {exc}"],
        }

    return {
        "file": path.name,
        "method": METHOD,
        "raw_text_chars": len(raw_markdown),
        "raw_markdown": raw_markdown,
        "warnings": warnings_list,
    }
