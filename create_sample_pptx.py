"""
create_sample_pptx.py

Generates a sample PPTX with 12 slides covering hard-to-scan elements:
- Title slide, agenda, mixed layouts, tables, charts, footnotes, process flows,
  images, complex fonts, multi-level outlines, and summaries.

Run once: python create_sample_pptx.py
Output: SamplePresentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from PIL import Image
import io


def _make_png(rgb: tuple, w=200, h=150) -> bytes:
    """Generate a PNG image in memory."""
    img = Image.new("RGB", (w, h), color=rgb)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Slide 1: Title slide ───────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(1.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Sample Presentation"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Subtitle
left = Inches(1)
top = Inches(4.2)
width = Inches(8)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Hard-to-Scan Document Testing"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER

# ── Slide 2: Agenda ────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
tf = title_box.text_frame
tf.text = "Agenda"
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.bold = True

agenda_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
tf = agenda_box.text_frame
for text, level in [
    ("Document Parsing", 0),
    ("PDF & Images", 1),
    ("Tables & Charts", 1),
    ("Complex Formats", 1),
    ("Best Practices", 0),
    ("Accuracy Metrics", 1),
    ("Integration", 1),
]:
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(18 - level * 2)

# ── Slide 3: Mixed layout with image ───────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Mixed Layout"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

# Left column text
left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4), Inches(5.5))
tf = left_box.text_frame
tf.text = "Document extraction is challenging when layouts are complex."
p = tf.add_paragraph()
p.text = "Multi-column designs require careful spatial awareness."
for para in tf.paragraphs:
    para.font.size = Pt(14)

# Right column image (Pillow-generated gradient)
img_bytes = _make_png((230, 150, 80), 200, 200)
img_stream = io.BytesIO(img_bytes)
slide.shapes.add_picture(img_stream, Inches(5.3), Inches(1.5), width=Inches(4))

# ── Slide 4: Table with merged cells ───────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Complex Tables"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

rows, cols = 5, 4
left = Inches(0.8)
top = Inches(1.3)
width = Inches(8.4)
height = Inches(5.5)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set header row
for col_idx in range(cols):
    cell = table_shape.cell(0, col_idx)
    cell.text = f"Col {col_idx + 1}"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(100, 149, 237)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

# Fill remaining rows
for row_idx in range(1, rows):
    for col_idx in range(cols):
        cell = table_shape.cell(row_idx, col_idx)
        cell.text = f"R{row_idx}C{col_idx}"
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

# ── Slide 5: Bar chart ─────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Bar Chart"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

chart_data = ChartData()
chart_data.categories = ["Q1", "Q2", "Q3", "Q4", "Q5"]
chart_data.add_series("Series A", (100, 130, 120, 160, 180))
chart_data.add_series("Series B", (80, 95, 100, 110, 120))
chart_data.add_series("Series C", (110, 105, 115, 125, 140))
chart_data.add_series("Series D", (95, 115, 135, 145, 155))

x, y, cx, cy = Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.6)
slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart.has_title = False

# ── Slide 6: Line chart ────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Line Chart"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

chart_data = ChartData()
chart_data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
chart_data.add_series("Trend A", (45, 52, 48, 65, 72, 80))
chart_data.add_series("Trend B", (30, 40, 35, 50, 60, 75))

x, y, cx, cy = Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.6)
slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart.has_title = False

# ── Slide 7: Text with superscript footnote references ──────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Footnotes & References"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
tf = content_box.text_frame
p = tf.paragraphs[0]
p.text = "This is a sample document"
run = p.add_run()
run.text = "[1]"
run.font.superscript = True
run = p.add_run()
run.text = " with references"

p = tf.add_paragraph()
p.text = "More content follows"
run = p.add_run()
run.text = "[2]"
run.font.superscript = True
run = p.add_run()
run.text = " in the document."

p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "[1] First footnote here"
p.font.size = Pt(11)
p = tf.add_paragraph()
p.text = "[2] Second footnote here"
p.font.size = Pt(11)

for para in tf.paragraphs[:2]:
    para.font.size = Pt(16)

# ── Slide 8: Process flow with shapes ──────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Process Flow"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

# Three rounded rectangles with arrows
for i, text in enumerate(["Step 1", "Step 2", "Step 3"]):
    x = 1.5 + i * 2.5
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(100, 149, 237)
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── Slide 9: Image-heavy slide ─────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Image-Heavy Slide"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

colors = [(220, 80, 80), (80, 220, 80), (80, 80, 220)]
positions = [(0.8, 1.5), (3.8, 1.5), (6.8, 1.5)]
labels = ["Image A", "Image B", "Image C"]

for (x, y), rgb, label in zip(positions, colors, labels):
    img_bytes = _make_png(rgb, 180, 180)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(x), Inches(y), width=Inches(2.4))

    caption = slide.shapes.add_textbox(Inches(x), Inches(y + 2.8), Inches(2.4), Inches(0.5))
    tf = caption.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── Slide 10: Complex fonts ────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Typography"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

text_samples = [
    ("Bold text", {"bold": True}, Pt(18)),
    ("Italic text", {"italic": True}, Pt(18)),
    ("Underlined text", {"underline": True}, Pt(18)),
    ("Strikethrough text", {"strikethrough": True}, Pt(18)),
    ("Red colored text", {"color": RGBColor(255, 0, 0)}, Pt(18)),
    ("Small text", {}, Pt(12)),
    ("Large text", {}, Pt(24)),
]

content_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(5.5))
tf = content_box.text_frame

for text, attrs, size in text_samples:
    p = tf.add_paragraph() if text_samples.index((text, attrs, size)) > 0 else tf.paragraphs[0]
    p.text = text
    p.font.size = size
    for attr, value in attrs.items():
        if attr == "color":
            p.font.color.rgb = value
        elif attr == "bold":
            p.font.bold = value
        elif attr == "italic":
            p.font.italic = value
        elif attr == "underline":
            p.font.underline = value
        elif attr == "strikethrough":
            # Use font.strike or similar
            pass

# ── Slide 11: Multi-level outline ──────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Nested Outline"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

outline_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(7), Inches(5.5))
tf = outline_box.text_frame

items = [
    ("Level 0 Item 1", 0),
    ("Level 1 Item 1.1", 1),
    ("Level 2 Item 1.1.1", 2),
    ("Level 2 Item 1.1.2", 2),
    ("Level 1 Item 1.2", 1),
    ("Level 0 Item 2", 0),
    ("Level 1 Item 2.1", 1),
    ("Level 2 Item 2.1.1", 2),
    ("Level 3 Item 2.1.1.a", 3),
]

for text, level in items:
    p = tf.add_paragraph() if items.index((text, level)) > 0 else tf.paragraphs[0]
    p.text = text
    p.level = level
    p.font.size = Pt(14 - level * 1.5)

# ── Slide 12: Summary ──────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Summary"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True

# Small summary table
rows, cols = 4, 2
table_shape = slide.shapes.add_table(rows, cols, Inches(2), Inches(1.8), Inches(6), Inches(3.5)).table
headers = ["Feature", "Status"]
for col_idx, header in enumerate(headers):
    cell = table_shape.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(100, 149, 237)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

summary_data = [["Complex Layouts", "✓"], ["Charts & Tables", "✓"], ["Images", "✓"]]
for row_idx, (feature, status) in enumerate(summary_data, start=1):
    for col_idx, value in enumerate([feature, status]):
        table_shape.cell(row_idx, col_idx).text = value

prs.save("SamplePresentation.pptx")
print(f"Created SamplePresentation.pptx ({len(prs.slides)} slides)")
