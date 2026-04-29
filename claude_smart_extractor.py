"""
claude_smart_extractor.py

Cost-optimized Claude extractor using liteparse's per-page trigger logic:

  Page needs Claude when ANY of:
    - native text length < 150 chars  (likely scanned)
    - page has embedded images        (charts, photos, diagrams)
    - text has high garbled-char ratio (corrupted text layer)
  Otherwise: free extraction via PyMuPDF (no API call).

For image files (jpg/png/…): always Claude.
For DOCX/PPTX/XLSX/text: same HTML→PDF conversion as claude_extractor, then per-page routing.

Requires: ANTHROPIC_API_KEY env var, PyMuPDF (fitz) installed.
"""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

METHOD = "claude_smart"

_IMAGE_EXTS  = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
_TEXT_EXTS   = {".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml", ".yaml", ".yml"}
_CONVERT_EXTS = {".docx", ".pptx", ".xlsx"} | _TEXT_EXTS

_EXTRACTION_PROMPT = """Extract all content from this document page and return it as Markdown.

Rules:
- Preserve the exact order content appears in the document — top to bottom, do not regroup.
- Use ## for section headings and ### for sub-headings, matching the document hierarchy.
- Render label:value pairs as **Label:** Value on their own line.
- Render tables as Markdown tables with headers and all rows.
- Copy all text verbatim — do not rephrase, summarize, or add commentary.
- Do not add any text that is not in the document.
- For charts, graphs, or plots that contain legible numeric or categorical data, render the
  data as a Markdown table (rows × columns matching the chart axes and legend), then
  immediately follow it with a one-line figure note describing the chart type, title, and
  any annotation not captured in the table:
  > **[Figure]** <chart type, title, axis labels, color encoding, key annotations>
- For every other visual element (logo, photograph, diagram, sketch, signature, illustration,
  stamp, or any chart whose data values cannot be read) insert only a figure block:
  > **[Figure]** <description of shape, colors, symbols, text, spatial layout, and meaning>
  Do not just name it — describe it as if explaining to someone who cannot see it.
"""

_MODEL      = "claude-sonnet-4-6"
_MAX_TOKENS = 16000   # per-page calls; 16K is enough for a single page

# --- trigger thresholds (mirrors liteparse) ---
_MIN_TEXT_LEN    = 150   # chars of native text below which we call Claude
_MAX_GARBLE_RATE = 0.05  # fraction of replacement/control chars above which text is garbled


def _is_garbled(text: str) -> bool:
    if not text:
        return False
    bad = sum(1 for c in text if c == "�" or (ord(c) < 32 and c not in "\n\r\t"))
    return (bad / len(text)) > _MAX_GARBLE_RATE


def _page_needs_claude(page) -> bool:
    """Mirror of liteparse processPageOcr trigger, using fitz page."""
    text     = page.get_text("text")
    text_len = len(text.strip())
    has_img  = len(page.get_images()) > 0
    return text_len < _MIN_TEXT_LEN or has_img or _is_garbled(text)


_MAX_IMAGE_BYTES = 4 * 1024 * 1024  # stay under Claude's 5 MB limit


def _render_page_png(page, dpi: int = 100) -> bytes:
    import fitz
    for d in (dpi, 72, 50):
        mat = fitz.Matrix(d / 72, d / 72)
        pix = page.get_pixmap(matrix=mat)
        data = pix.tobytes("png")
        if len(data) <= _MAX_IMAGE_BYTES:
            return data
    # last resort: render at 50 DPI and accept it
    return data


def _extract_page_native(page) -> str:
    try:
        return page.get_text("markdown")
    except Exception:
        return page.get_text("text")


def _call_page_image(client, png_bytes: bytes, page_num: int, file_name: str) -> str:
    with client.messages.stream(
        model=_MODEL,
        max_tokens=_MAX_TOKENS,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/png",
                        "data": base64.standard_b64encode(png_bytes).decode(),
                    },
                },
                {"type": "text", "text": f"{_EXTRACTION_PROMPT}\n\nDOCUMENT: {file_name}, page {page_num}"},
            ],
        }],
    ) as stream:
        return stream.get_final_text()


def _call_with_image_file(client, content: bytes, ext: str, file_name: str) -> str:
    media_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                   ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}
    with client.messages.stream(
        model=_MODEL,
        max_tokens=_MAX_TOKENS,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_types.get(ext, "image/jpeg"),
                        "data": base64.standard_b64encode(content).decode(),
                    },
                },
                {"type": "text", "text": f"{_EXTRACTION_PROMPT}\n\nDOCUMENT: {file_name}"},
            ],
        }],
    ) as stream:
        return stream.get_final_text()


def _extract_pdf_smart(client, pdf_bytes: bytes, file_name: str) -> tuple[str, list[str]]:
    """Route each page to native extraction or Claude based on trigger conditions."""
    import fitz

    doc    = fitz.open(stream=pdf_bytes, filetype="pdf")
    parts  = []
    warnings: list[str] = []
    claude_pages: list[int] = []
    native_pages: list[int] = []

    for page_num, page in enumerate(doc, start=1):
        if _page_needs_claude(page):
            claude_pages.append(page_num)
            try:
                png = _render_page_png(page)
                md  = _call_page_image(client, png, page_num, file_name)
            except Exception as exc:
                md = _extract_page_native(page)
                warnings.append(f"Page {page_num}: Claude failed ({exc}), fell back to native")
        else:
            native_pages.append(page_num)
            md = _extract_page_native(page)
        parts.append(md)

    if claude_pages:
        warnings.append(f"Claude (vision) used for pages: {claude_pages}")
    if native_pages:
        warnings.append(f"Native PyMuPDF used for pages: {native_pages}")

    return "\n\n---\n\n".join(parts), warnings


# --- format converters (copied from claude_extractor) ---

def _sniff_image_media_type(blob: bytes) -> str:
    if blob[:2] == b"\xff\xd8":     return "image/jpeg"
    if blob[:4] == b"\x89PNG":      return "image/png"
    if blob[:4] in (b"GIF8", b"GIF9"): return "image/gif"
    return "image/png"


def _html_to_pdf_bytes(html: str) -> bytes:
    import fitz
    story    = fitz.Story(html)
    mediabox = fitz.paper_rect("a4")
    where    = mediabox + (50, 50, -50, -50)
    buf      = io.BytesIO()
    writer   = fitz.DocumentWriter(buf)
    more = True
    while more:
        dev  = writer.begin_page(mediabox)
        more, _ = story.place(where)
        story.draw(dev)
        writer.end_page()
    writer.close()
    return buf.getvalue()


def _docx_to_html(content: bytes) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    doc     = Document(io.BytesIO(content))
    img_map: dict[str, str] = {}
    for rId, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            mt   = _sniff_image_media_type(blob)
            img_map[rId] = f"data:{mt};base64,{base64.b64encode(blob).decode()}"
    parts = ["<html><body>"]
    for block in doc.element.body:
        tag = block.tag.split("}")[-1]
        if tag == "p":
            style = getattr(getattr(block, "style", None), "name", "") or ""
            text  = block.text.strip() if hasattr(block, "text") else ""
            if not text:
                parts.append("<br/>"); continue
            if   "Heading 1" in style: parts.append(f"<h1>{text}</h1>")
            elif "Heading 2" in style: parts.append(f"<h2>{text}</h2>")
            elif "Heading 3" in style: parts.append(f"<h3>{text}</h3>")
            elif "List"      in style: parts.append(f"<li>{text}</li>")
            else:                      parts.append(f"<p>{text}</p>")
            for blip in block.iter(qn("a:blip")):
                embed = blip.get(qn("r:embed"))
                if embed and embed in img_map:
                    parts.append(f'<img src="{img_map[embed]}" style="max-width:100%;"/>')
        elif tag == "tbl":
            parts.append("<table border='1' cellpadding='4'>")
            for row in block:
                if row.tag.split("}")[-1] != "tr": continue
                parts.append("<tr>")
                for cell in row:
                    if cell.tag.split("}")[-1] != "tc": continue
                    parts.append(f"<td>{''.join(p.text or '' for p in cell.iter() if p.tag.endswith('}t'))}</td>")
                parts.append("</tr>")
            parts.append("</table>")
    parts.append("</body></html>")
    return "\n".join(parts)


def _pptx_to_html(content: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs   = Presentation(io.BytesIO(content))
    parts = ["<html><body>"]
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"<h2>Slide {i}</h2>")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip(): parts.append(f"<p>{para.text.strip()}</p>")
            if shape.has_table:
                parts.append("<table border='1' cellpadding='4'>")
                for row in shape.table.rows:
                    parts.append("<tr>" + "".join(f"<td>{c.text.strip()}</td>" for c in row.cells) + "</tr>")
                parts.append("</table>")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    b = shape.image.blob
                    parts.append(f'<img src="data:{_sniff_image_media_type(b)};base64,{base64.b64encode(b).decode()}" style="max-width:100%;"/>')
                except Exception: pass
        parts.append("<hr/>")
    parts.append("</body></html>")
    return "\n".join(parts)


def _xlsx_to_html(content: bytes) -> str:
    import openpyxl
    wb    = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = ["<html><body>"]
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"<h2>{name}</h2><table border='1' cellpadding='4'>")
        for row in ws.iter_rows(values_only=True):
            parts.append("<tr>" + "".join(f"<td>{v or ''}</td>" for v in row) + "</tr>")
        parts.append("</table>")
    parts.append("</body></html>")
    return "\n".join(parts)


# --- public entry point ---

def extract(file_path: str) -> dict:
    path    = Path(file_path)
    ext     = path.suffix.lower()
    content = path.read_bytes()
    warnings_list: list[str] = []

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {"file": path.name, "method": METHOD, "raw_markdown": "",
                "raw_text_chars": 0, "warnings": ["ANTHROPIC_API_KEY not set — skipping"]}

    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
    except ImportError:
        return {"file": path.name, "method": METHOD, "raw_markdown": "",
                "raw_text_chars": 0, "warnings": ["anthropic package not installed"]}

    try:
        if ext in _IMAGE_EXTS:
            # Images have no native text layer — always Claude
            raw_markdown = _call_with_image_file(client, content, ext, path.name)

        elif ext == ".pdf":
            raw_markdown, warnings_list = _extract_pdf_smart(client, content, path.name)

        elif ext in _CONVERT_EXTS:
            if   ext == ".docx": html = _docx_to_html(content)
            elif ext == ".pptx": html = _pptx_to_html(content)
            elif ext == ".xlsx": html = _xlsx_to_html(content)
            else:
                text = content.decode("utf-8", errors="replace")
                html = f"<html><body><pre>{text}</pre></body></html>"
            pdf_bytes    = _html_to_pdf_bytes(html)
            raw_markdown, warnings_list = _extract_pdf_smart(client, pdf_bytes, path.name)

        else:
            try:
                text         = content.decode("utf-8", errors="replace")
                html         = f"<html><body><pre>{text}</pre></body></html>"
                pdf_bytes    = _html_to_pdf_bytes(html)
                raw_markdown, warnings_list = _extract_pdf_smart(client, pdf_bytes, path.name)
                warnings_list.append(f"Unknown type {ext}; treated as plain text")
            except Exception:
                return {"file": path.name, "method": METHOD, "raw_markdown": "",
                        "raw_text_chars": 0, "warnings": [f"Unsupported file type: {ext}"]}

    except Exception as exc:
        return {"file": path.name, "method": METHOD, "raw_markdown": "",
                "raw_text_chars": 0, "warnings": [f"Error: {exc}"]}

    return {
        "file":           path.name,
        "method":         METHOD,
        "raw_text_chars": len(raw_markdown),
        "raw_markdown":   raw_markdown,
        "warnings":       warnings_list,
    }
